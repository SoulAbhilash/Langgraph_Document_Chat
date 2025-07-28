# --- Standard Library Imports --- #
import io
from typing import List

# --- Third-party Imports --- #
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from streamlit.runtime.uploaded_file_manager import UploadedFile
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter


class DocumentsHandler:
    """
    Handles text extraction and chunking from uploaded documents.

    Supported file types:
        - PDF
        - DOCX (Word)
        - PPTX (PowerPoint)
        - XLSX (Excel)
    """

    def __init__(self, files: List[UploadedFile]) -> None:
        """
        Initializes the handler with uploaded files.

        Args:
            files (List[UploadedFile]): List of Streamlit-uploaded files.
        """
        self._files: List[UploadedFile] = files or []

    def create_chunks(self) -> List[Document]:
        """
        Splits parsed documents into manageable text chunks.

        Returns:
            List[Document]: Chunked documents with metadata.
        """
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        to_documents = self._create_documents_list()
        chunks = splitter.split_documents(to_documents)
        return chunks

    def _create_documents_list(self) -> List[Document]:
        """
        Parses all uploaded documents into a list of LangChain Document objects.

        Returns:
            List[Document]: Parsed documents.
        """
        all_documents: List[Document] = []

        for file in self._files:
            filetype = file.type

            if filetype == "application/pdf":
                docs = self._extract_text_from_pdf(file)
            elif (
                filetype
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                docs = self._extract_text_from_word(file)
            elif (
                filetype
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                docs = self._extract_text_from_ppt(file)
            elif (
                filetype
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ):
                docs = self._extract_text_from_excel(file)
            else:
                continue  # Skip unsupported file types

            all_documents.extend(docs)

        return all_documents

    def _extract_text_from_pdf(self, file: UploadedFile) -> List[Document]:
        """
        Extracts text from each page of a PDF.

        Args:
            file (UploadedFile): The uploaded PDF file.

        Returns:
            List[Document]: Text content from PDF pages with metadata.
        """
        pdf_reader = PdfReader(file)
        to_documents: List[Document] = []

        for idx, page in enumerate(pdf_reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                to_documents.append(
                    Document(
                        page_content=text,
                        metadata={"page": idx, "filename": file.name, "type": "pdf"},
                    )
                )

        return to_documents

    def _extract_text_from_word(self, file: UploadedFile) -> List[Document]:
        """
        Extracts and combines text from a DOCX (Word) file.

        Args:
            file (UploadedFile): The uploaded Word document.

        Returns:
            List[Document]: One document with all paragraph text.
        """
        doc = DocxDocument(file)
        full_text = "\n".join(
            [para.text for para in doc.paragraphs if para.text.strip()]
        )

        return [
            Document(
                page_content=full_text, metadata={"filename": file.name, "type": "word"}
            )
        ]

    def _extract_text_from_ppt(self, file: UploadedFile) -> List[Document]:
        """
        Extracts text from each slide in a PowerPoint file.

        Args:
            file (UploadedFile): The uploaded PPTX file.

        Returns:
            List[Document]: One document per slide with metadata.
        """
        presentation = Presentation(file)
        to_documents: List[Document] = []

        for idx, slide in enumerate(presentation.slides):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            slide_text = "\n".join(texts)
            if slide_text:
                to_documents.append(
                    Document(
                        page_content=slide_text,
                        metadata={"slide": idx, "filename": file.name, "type": "ppt"},
                    )
                )

        return to_documents

    def _extract_text_from_excel(self, file: UploadedFile) -> List[Document]:
        """
        Extracts text from all sheets in an Excel file.

        Args:
            file (UploadedFile): The uploaded XLSX file.

        Returns:
            List[Document]: One document per sheet with metadata.
        """
        to_documents: List[Document] = []
        in_memory_file = io.BytesIO(file.read())
        xls = pd.ExcelFile(in_memory_file)

        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            text = df.to_string(index=False)
            if text.strip():
                to_documents.append(
                    Document(
                        page_content=text,
                        metadata={
                            "sheet": sheet_name,
                            "filename": file.name,
                            "type": "excel",
                        },
                    )
                )

        return to_documents


if __name__ == "__main__":
    # Simple test run (for debugging only)
    handler = DocumentsHandler([])
    print("Initialized with no files:", handler)
